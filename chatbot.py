import streamlit as st
from anthropic import Anthropic
from dotenv import load_dotenv
import base64
import fitz
from pptx import Presentation
from docx import Document
import io
import os
import re

load_dotenv()
client = Anthropic()

with open("prompt.txt", "r", encoding="utf-8") as f:
    SYSTEM_PROMPT = f.read()

def _get_secret(key):
    """secrets.toml(클라우드) 또는 환경변수(.env) 어느 쪽이든 안전하게 읽기"""
    try:
        return st.secrets[key]
    except Exception:
        return os.environ.get(key)

def check_password():
    """APP_PASSWORD가 설정돼 있으면 비밀번호 잠금. 미설정(로컬)이면 통과."""
    app_pw = _get_secret("APP_PASSWORD")
    if not app_pw:
        return True
    if st.session_state.get("authenticated"):
        return True

    def _verify():
        if st.session_state.get("pw_input") == app_pw:
            st.session_state["authenticated"] = True
            st.session_state.pop("pw_input", None)
        else:
            st.session_state["authenticated"] = False

    st.markdown("### RD부 업무 안내 챗봇")
    st.text_input("비밀번호를 입력하세요", type="password",
                  key="pw_input", on_change=_verify)
    if st.session_state.get("authenticated") is False:
        st.error("비밀번호가 올바르지 않습니다.")
    return False

def get_image_base64(image_path):
    try:
        with open(image_path, "rb") as f:
            return base64.b64encode(f.read()).decode()
    except:
        return None

def extract_pptx_text(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for i, slide in enumerate(prs.slides):
        text += f"\n[슬라이드 {i+1}]\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return text

# ── 선택지시문 자동 점검용 패턴 ─────────────────────────────
# 선택지시문: PROG/BASE 대괄호는 제외하고, 실제 지시문 키워드가 든 [...]만 인정
_INSTR_RE   = re.compile(r'\[(?![^\]]*(?:PROG|BASE))[^\]]*(?:선택|직접\s*입력|수치형|순위|척도)[^\]]*\]')
_LABEL_RE   = re.compile(r'^\s*(SQ|DQ|Q)\s*\d+[A-Z]?(?:-\d+)?')
_CHOICE_RE  = re.compile(r'^\s*\d+\s*\)')
_PROGBASE_RE = re.compile(r'^\s*\[(PROG|BASE)\b')


def _is_question_start(t):
    """문단이 새로운 '문항의 시작'인지 판정. 보기줄/PROG·BASE줄/지시문 연속줄은 제외."""
    t = t.strip()
    if not t or _PROGBASE_RE.match(t):
        return False
    if _LABEL_RE.match(t):          # SQ7, Q19, DQ1 등 명시 라벨
        return True
    if _CHOICE_RE.match(t):         # "1) ..." 보기줄
        return False
    if '?' in t:                    # 물음표가 있으면 질문 문장
        return True
    return False


def _question_id(t):
    t = t.strip()
    m = _LABEL_RE.match(t)
    if m:
        return re.match(r'^\s*((SQ|DQ|Q)\s*\d+[A-Z]?(?:-\d+)?)', t).group(1).replace(' ', '')
    return t[:24].strip()           # 라벨 없으면 앞부분으로 식별


def _audit_choice_instructions(para_texts):
    """문항을 블록 단위로 묶어, 블록 전체에 선택지시문이 하나도 없는 문항만 골라낸다.
    (선택지시문이 질문 문장과 다른 줄에 분리돼 있어도 같은 블록이면 '있음'으로 인정)"""
    blocks = []
    cur = None
    for txt in para_texts:
        if _is_question_start(txt):
            cur = {'id': _question_id(txt), 'has': bool(_INSTR_RE.search(txt))}
            blocks.append(cur)
        elif cur is not None:
            if _INSTR_RE.search(txt):
                cur['has'] = True
        # cur가 없을 때(머리말 등)는 문항이 아니므로 버린다
    missing = [b['id'] for b in blocks if not b['has']]
    return len(blocks), missing


def extract_docx_text(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    text = ""
    ns = 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
    para_texts = []  # 선택지시문 자동 점검용 (본문 문단 텍스트만 수집)
    num_counts = {}  # {(numId, ilvl): count} — 자동번호매기기 카운터

    for child in doc.element.body:
        tag = child.tag

        if tag == f'{{{ns}}}p':
            para_text = ''.join(node.text or '' for node in child.iter()
                                if node.tag == f'{{{ns}}}t')
            if para_text.strip():
                para_texts.append(para_text.replace('\xa0', ' '))

            # 자동번호매기기(numPr) 감지 후 번호 복원
            num_prefix = ''
            pPr = child.find(f'{{{ns}}}pPr')
            if pPr is not None:
                numPr = pPr.find(f'{{{ns}}}numPr')
                if numPr is not None:
                    ilvl_elem = numPr.find(f'{{{ns}}}ilvl')
                    numId_elem = numPr.find(f'{{{ns}}}numId')
                    if ilvl_elem is not None and numId_elem is not None:
                        ilvl = int(ilvl_elem.get(f'{{{ns}}}val', 0))
                        numId = int(numId_elem.get(f'{{{ns}}}val', 0))
                        if numId != 0:
                            key = (numId, ilvl)
                            num_counts[key] = num_counts.get(key, 0) + 1
                            num_prefix = f'{num_counts[key]}) '

            if para_text.strip():
                text += num_prefix + para_text + "\n"

        elif tag == f'{{{ns}}}tbl':
            rows = [r for r in child if r.tag == f'{{{ns}}}tr']
            num_rows = len(rows)
            if num_rows > 1:
                attr_count = num_rows - 1
                first_row_cells = [tc for tc in rows[0] if tc.tag == f'{{{ns}}}tc']
                choice_count = max(0, len(first_row_cells) - 1)
                text += f"[테이블: 속성 {attr_count}개, 보기 {choice_count}개]\n"
                for i, row in enumerate(rows):
                    cells = [tc for tc in row if tc.tag == f'{{{ns}}}tc']
                    cell_texts = []
                    for tc in cells:
                        t_text = ''.join(node.text or '' for node in tc.iter()
                                         if node.tag == f'{{{ns}}}t')
                        cell_texts.append(t_text.strip())
                    label = "[헤더]" if i == 0 else f"[속성{i}]"
                    text += f"  {label} {' | '.join(cell_texts)}\n"

    # ── 선택지시문 자동 점검 결과 부착 (LLM이 눈으로 세지 말고 이 결과를 신뢰) ──
    q_count, missing = _audit_choice_instructions(para_texts)
    text += "\n\n[선택지시문 자동 점검 결과]\n"
    text += f"- 점검한 문항 블록 수: {q_count}\n"
    if missing:
        text += "- 선택지시문이 없는 문항: " + ", ".join(missing) + "\n"
    else:
        text += "- 선택지시문이 없는 문항: 없음 (모든 문항에 선택지시문 있음)\n"

    return text

def extract_pdf_text(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = ""
    for i, page in enumerate(doc):
        text += f"\n[페이지 {i+1}]\n"
        text += page.get_text()
    return text

def stream_ai_response(api_messages):
    with client.messages.stream(
        model="claude-sonnet-4-6",
        max_tokens=16384,   # 긴 설문지 검수 결과가 잘리지 않도록 (사용한 만큼만 과금)
        system=[{
            "type": "text",
            "text": SYSTEM_PROMPT,
            "cache_control": {"type": "ephemeral"}
        }],
        messages=api_messages
    ) as stream:
        buffer = ""
        for text in stream.text_stream:
            buffer += text
            if "\n\n" in buffer:
                parts = buffer.split("\n\n")
                for part in parts[:-1]:
                    yield part + "\n\n"
                buffer = parts[-1]
        if buffer:
            yield buffer

st.set_page_config(
    page_title="RD부 업무 안내 챗봇",
    page_icon="로고.png",
    layout="centered"
)

# 라이트/다크 팔레트 (CSS 변수로 주입)
LIGHT_THEME = {
    "bg": "#FAF9F5",
    "panel": "#F5F4EE",
    "border": "#E8E6DC",
    "border2": "#E0DDD0",
    "text": "#1F1E1D",
    "muted": "#8A8678",
    "muted2": "#9B9786",
    "bubble": "#F0EEE6",
    "input-bg": "#FFFFFF",
    "chip-bg": "#FFFFFF",
    "chip-hover": "#F0EEE6",
    "chip-border-hover": "#C8C5B5",
    "warn-bg": "#F7F0EA",
    "warn-border": "#E8DBCC",
    "warn-text": "#6B5D4F",
    "warn-title": "#A0633A",
    "shadow": "rgba(31,30,29,0.05)",
}
DARK_THEME = {
    "bg": "#262624",
    "panel": "#1F1E1D",
    "border": "#3A3A37",
    "border2": "#44443F",
    "text": "#ECECEA",
    "muted": "#A6A296",
    "muted2": "#8F8B80",
    "bubble": "#3A3A37",
    "input-bg": "#30302E",
    "chip-bg": "#30302E",
    "chip-hover": "#3A3A37",
    "chip-border-hover": "#55554F",
    "warn-bg": "#332B24",
    "warn-border": "#4A3D30",
    "warn-text": "#C9B49C",
    "warn-title": "#D08D5E",
    "shadow": "rgba(0,0,0,0.25)",
}

_palette = DARK_THEME if st.session_state.get("dark_mode") else LIGHT_THEME
_css_vars = "".join(f"--{k}: {v};" for k, v in _palette.items())
st.markdown(f"<style>:root {{{_css_vars}}}</style>", unsafe_allow_html=True)

st.markdown("""
<style>
@import url('https://cdn.jsdelivr.net/gh/orioncactus/pretendard@v1.3.9/dist/web/static/pretendard.min.css');
@import url('https://fonts.googleapis.com/css2?family=Noto+Serif+KR:wght@500;600&display=swap');

html, body, [class*="st-"], .stApp {
    font-family: 'Pretendard', 'Noto Sans KR', sans-serif;
}

/* Material 아이콘 폰트는 덮어쓰지 않기 (사이드바 접기 버튼 등) */
[data-testid="stIconMaterial"],
span.material-symbols-rounded,
[data-testid="stSidebarCollapseButton"] span,
[data-testid="stExpandSidebarButton"] span {
    font-family: 'Material Symbols Rounded' !important;
}

.stApp { background-color: var(--bg); }
#MainMenu, footer, [data-testid="stToolbar"] { visibility: hidden; }
[data-testid="stHeader"] { background: transparent; }

/* 사이드바 다시 열기(>>) 버튼은 stToolbar 안에 있어 위 규칙에 같이 숨겨짐 → 복구 */
[data-testid="stExpandSidebarButton"] {
    visibility: visible !important;
    color: var(--muted) !important;
}

/* 본문 텍스트 (다크모드 대응) */
[data-testid="stMarkdownContainer"],
[data-testid="stMarkdownContainer"] p,
[data-testid="stMarkdownContainer"] li,
[data-testid="stMarkdownContainer"] h1,
[data-testid="stMarkdownContainer"] h2,
[data-testid="stMarkdownContainer"] h3,
[data-testid="stMarkdownContainer"] h4,
[data-testid="stMarkdownContainer"] strong,
[data-testid="stMarkdownContainer"] td,
[data-testid="stMarkdownContainer"] th {
    color: var(--text);
}

/* ── 다크모드 토글 버튼 (우상단 고정) ── */
.st-key-theme_toggle {
    position: fixed;
    top: 10px;
    right: 24px;
    z-index: 1000000;  /* Streamlit 헤더(999990)보다 위 */
    width: auto !important;
}
.st-key-theme_toggle button {
    background: transparent !important;
    border: 1px solid var(--border2) !important;
    color: var(--muted) !important;
    border-radius: 999px !important;
    font-size: 12.5px !important;
    font-weight: 500 !important;
    min-height: 32px !important;
    padding: 2px 14px !important;
    box-shadow: none !important;
}
.st-key-theme_toggle button:hover {
    background: var(--chip-hover) !important;
    color: var(--text) !important;
}

/* ── 사이드바 ── */
[data-testid="stSidebar"] {
    background: var(--panel) !important;
    border-right: 1px solid var(--border);
}
[data-testid="stSidebar"] * { color: var(--text); }

.side-section { margin-bottom: 4px; }
.side-title {
    font-size: 11.5px;
    font-weight: 600;
    letter-spacing: 0.6px;
    color: var(--muted);
    margin: 22px 0 10px;
}
.side-item {
    font-size: 13.5px;
    color: var(--text);
    line-height: 1.5;
    padding: 2px 0;
}
.side-sub {
    font-size: 12px;
    color: var(--muted2);
    padding-left: 12px;
    line-height: 1.5;
}
.side-divider { border-top: 1px solid var(--border); margin: 18px 0 4px; }
.side-warning {
    background: var(--warn-bg);
    border: 1px solid var(--warn-border);
    border-radius: 10px;
    padding: 12px 14px;
    font-size: 12.5px;
    line-height: 1.7;
    color: var(--warn-text);
    margin-top: 10px;
}
.side-warning, .side-warning span { color: var(--warn-text) !important; }
.side-warning .warn-title {
    font-weight: 600;
    font-size: 12px;
    color: var(--warn-title) !important;
    margin-bottom: 4px;
}
[data-testid="stSidebar"] .stButton button {
    background: transparent !important;
    border: 1px solid var(--border2) !important;
    color: var(--text) !important;
    border-radius: 10px !important;
    font-size: 13.5px !important;
    font-weight: 500 !important;
    box-shadow: none !important;
}
[data-testid="stSidebar"] .stButton button:hover {
    background: var(--chip-hover) !important;
    border-color: var(--chip-border-hover) !important;
}

/* ── 사용자 메시지: 오른쪽 연한 박스 ── */
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"]) {
    display: flex !important;
    flex-direction: row !important;
    justify-content: flex-end !important;
    background: transparent !important;
    box-shadow: none !important;
    padding: 8px 0 !important;
    gap: 0 !important;
}
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"])
[data-testid="stChatMessageContent"] {
    background: var(--bubble) !important;
    color: var(--text) !important;
    border-radius: 16px !important;
    box-shadow: none !important;
    max-width: 75% !important;
    width: fit-content !important;
    padding: 12px 18px !important;
    margin: 0 !important;
    line-height: 1.65 !important;
    display: inline-block !important;
}
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"])
[data-testid="stChatMessageContent"] p {
    color: var(--text) !important;
    margin: 0 !important;
}
/* Streamlit이 stMarkdownContainer에 margin-bottom:-1rem을 줘서 (마지막 p의 +1rem과 상쇄용)
   p margin을 0으로 만들면 -1rem만 남아 말풍선 높이가 16px 모자라 텍스트가 아래로 쏠림 → 상쇄 제거 */
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"])
[data-testid="stChatMessageContent"] [data-testid="stMarkdownContainer"] {
    margin-bottom: 0 !important;
}
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"])
[data-testid="stChatMessageAvatarUser"] {
    display: none !important;
}

/* ── 챗봇 메시지: 말풍선 없이 본문 텍스트 ── */
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarAssistant"]) {
    background: transparent !important;
    box-shadow: none !important;
    padding: 8px 0 !important;
    align-items: flex-start !important;
}
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarAssistant"])
[data-testid="stChatMessageContent"] {
    background: transparent !important;
    color: var(--text) !important;
    border-radius: 0 !important;
    box-shadow: none !important;
    max-width: 100% !important;
    padding: 0 !important;
    margin: 0 !important;
    line-height: 1.75 !important;
}
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarAssistant"])
[data-testid="stChatMessageAvatarAssistant"] {
    display: none !important;
}

/* ── 입력창 ── */
[data-testid="stChatInput"] {
    background: var(--input-bg) !important;
    border: 1px solid var(--border2) !important;
    border-radius: 18px !important;
    box-shadow: 0 2px 10px var(--shadow) !important;
}
[data-testid="stChatInput"] > div,
[data-testid="stChatInput"] div[data-baseweb="textarea"],
[data-testid="stChatInput"] div[data-baseweb="base-input"] {
    background: var(--input-bg) !important;
    border-color: transparent !important;
}
[data-testid="stChatInput"] textarea {
    background: transparent !important;
    font-size: 14.5px !important;
    color: var(--text) !important;
}
[data-testid="stChatInput"] textarea::placeholder {
    color: var(--muted) !important;
}
[data-testid="stChatInput"] button {
    color: var(--muted) !important;
}
[data-testid="stBottom"],
[data-testid="stBottom"] > div {
    background: var(--bg) !important;
}

/* ── 추천 질문 칩 ── */
.stButton button {
    border-radius: 12px !important;
    border: 1px solid var(--border2) !important;
    background: var(--chip-bg) !important;
    color: var(--text) !important;
    font-size: 13.5px !important;
    font-weight: 400 !important;
    box-shadow: none !important;
    transition: background 0.15s, border-color 0.15s !important;
}
.stButton button:hover {
    background: var(--chip-hover) !important;
    border-color: var(--chip-border-hover) !important;
    color: var(--text) !important;
}

/* ── 첫 화면 인사 ── */
.welcome {
    text-align: center;
    padding: 90px 0 36px;
}
.welcome-title {
    font-family: 'Noto Serif KR', serif;
    font-size: 29px;
    font-weight: 600;
    color: var(--text);
    letter-spacing: -0.3px;
}
.welcome-sub {
    font-size: 14.5px;
    color: var(--muted);
    margin-top: 12px;
    line-height: 1.6;
}

/* ── 파일 분석 중 표시 ── */
@keyframes gentle-pulse {
    0%, 100% { opacity: 0.45; }
    50% { opacity: 1; }
}
.analyzing {
    color: var(--muted);
    font-size: 13.5px;
    animation: gentle-pulse 1.6s ease-in-out infinite;
}
</style>
""", unsafe_allow_html=True)

# 비밀번호 잠금 (APP_PASSWORD 설정 시에만 동작)
if not check_password():
    st.stop()

# 상태 초기화
if "messages" not in st.session_state:
    st.session_state.messages = []
if "pending_question" not in st.session_state:
    st.session_state.pending_question = None
if "dark_mode" not in st.session_state:
    st.session_state.dark_mode = False

# 다크모드 토글 (우상단 고정 — CSS .st-key-theme_toggle)
_toggle_label = "라이트 모드" if st.session_state.dark_mode else "다크 모드"
if st.button(_toggle_label, key="theme_toggle"):
    st.session_state.dark_mode = not st.session_state.dark_mode
    st.rerun()

SAMPLE_QUESTIONS = [
    "이 설문지는 어느 부서로 의뢰 해야 하나요?",
    "이 조사의 적립금은 대략 얼마정도 인가요?",
    "이지서베이 조사는 어디로 의뢰하나요?",
    "웹업 프로세스가 어떻게 되나요?",
]

# 사이드바
with st.sidebar:
    logo_b64 = get_image_base64("로고.png")
    if logo_b64:
        st.markdown(f"""
            <div style="background:white;border:1px solid var(--border);
                        border-radius:12px;padding:14px;
                        text-align:center;margin-bottom:8px;">
                <img src="data:image/png;base64,{logo_b64}"
                     style="width:100%;max-width:160px;">
            </div>
        """, unsafe_allow_html=True)

    st.markdown("""
    <div class="side-title">문의 가능 내용</div>
    <div class="side-item">조사 의뢰처 안내</div>
    <div class="side-sub">RD부 · 온라인실사팀 결정</div>
    <div class="side-item">설문지 검토</div>
    <div class="side-sub">로직 · 오류 등 이상 여부 확인</div>
    <div class="side-item">예상 적립금 안내</div>
    <div class="side-sub">응답시간 기반 계산</div>
    <div class="side-item">업무 일정 안내</div>
    <div class="side-sub">웹업 · 실사 · 테이블 일정</div>

    <div class="side-title" style="margin-top:26px;">첨부 가능 파일</div>
    <div class="side-item">이미지(png · jpg) / PDF / PPT / 워드</div>
    <div class="side-sub">HWP는 PDF로 변환 후 첨부</div>
    <div class="side-sub">파일은 채팅창에 끌어다 놓으면 첨부됩니다</div>

    <div class="side-warning">
        <div class="warn-title">보안 주의사항</div>
        고객사명 · 기밀 프로젝트, 응답자 개인정보,
        사내 기밀 문서는 입력하지 마세요.<br>
        <span style="font-size:11.5px;opacity:0.8;">
        입력 내용은 외부 AI 서버를 경유합니다.</span>
    </div>
    <div class="side-divider"></div>
    """, unsafe_allow_html=True)

    if st.button("새 대화 시작", use_container_width=True):
        st.session_state.messages = []
        st.session_state.pending_question = None
        st.rerun()

# 첫 화면 인사 + 추천 질문
if not st.session_state.messages:
    st.markdown("""
    <div class="welcome">
        <div class="welcome-title">안녕하세요, RD부 업무 안내 챗봇입니다</div>
        <div class="welcome-sub">설문조사 의뢰처, ISAS 작업 가능 여부, 예상 적립금 등을 문의해 보세요.<br>
        설문지 파일은 아래 채팅창에 끌어다 놓으면 첨부됩니다.</div>
    </div>
    """, unsafe_allow_html=True)
    cols = st.columns(2)
    for i, q in enumerate(SAMPLE_QUESTIONS):
        with cols[i % 2]:
            if st.button(q, key=f"sq_{i}", use_container_width=True):
                st.session_state.pending_question = q
                st.rerun()

# 이전 대화 표시
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.markdown(msg["content"])

# 채팅 입력 (파일 첨부 일체형 — 드래그 앤 드롭 지원)
chat_value = st.chat_input(
    "문의 내용을 입력하거나 파일을 끌어다 놓으세요",
    accept_file=True,
    file_type=["png", "jpg", "jpeg", "pdf", "pptx", "docx"],
)

user_input = None
uploaded_file = None
if chat_value:
    user_input = (chat_value.text or "").strip()
    if chat_value.files:
        uploaded_file = chat_value.files[0]
        if not user_input:
            user_input = "첨부한 파일을 검토해 주세요."

if st.session_state.pending_question:
    user_input = st.session_state.pending_question
    st.session_state.pending_question = None

# 메시지 처리
if user_input:
    # ① 사용자 말풍선을 즉시 화면에 표시 (파일 파싱·답변 생성 전 — 체감 지연 제거)
    if uploaded_file:
        _fname = uploaded_file.name.lower()
        _tag = "이미지" if _fname.endswith((".png", ".jpg", ".jpeg")) else "파일"
        display_content = f"[{_tag} 첨부: {uploaded_file.name}]\n{user_input}"
    else:
        display_content = user_input
    with st.chat_message("user"):
        st.markdown(display_content)

    # ② 파일 파싱
    file_text = ""
    is_image = False
    image_data = None
    image_type = None

    if uploaded_file:
        file_bytes = uploaded_file.read()
        file_name = uploaded_file.name.lower()
        if file_name.endswith((".png", ".jpg", ".jpeg")):
            is_image = True
            image_data = base64.b64encode(file_bytes).decode()
            image_type = uploaded_file.type
        elif file_name.endswith(".pdf"):
            file_text = extract_pdf_text(file_bytes)
        elif file_name.endswith(".pptx"):
            file_text = extract_pptx_text(file_bytes)
        elif file_name.endswith(".docx"):
            file_text = extract_docx_text(file_bytes)

    if is_image and image_data:
        api_message = {
            "role": "user",
            "content": [
                {"type": "image", "source": {
                    "type": "base64",
                    "media_type": image_type,
                    "data": image_data
                }},
                {"type": "text",
                 "text": f"첨부 이미지를 참고해서 답변해줘.\n\n{user_input}"}
            ]
        }
    elif file_text:
        api_message = {
            "role": "user",
            "content": f"첨부 파일 내용:\n{file_text}\n\n"
                       f"위 내용을 참고해서 답변해줘.\n\n{user_input}"
        }
    else:
        api_message = {"role": "user", "content": user_input}

    st.session_state.messages.append({
        "role": "user", "content": display_content
    })

    api_messages = []
    for msg in st.session_state.messages[:-1]:
        api_messages.append({
            "role": msg["role"], "content": msg["content"]
        })
    api_messages.append(api_message)

    with st.chat_message("assistant"):
        has_file = bool(file_text or is_image)
        status_placeholder = st.empty()
        status_text = ("첨부 파일을 분석하고 있습니다" if has_file
                       else "답변을 준비하고 있습니다")
        status_placeholder.markdown(
            f'<div class="analyzing">{status_text}</div>',
            unsafe_allow_html=True
        )

        def _stream_clearing_status():
            first = True
            for chunk in stream_ai_response(api_messages):
                if first:
                    status_placeholder.empty()
                    first = False
                yield chunk

        answer = st.write_stream(_stream_clearing_status())

    st.session_state.messages.append({
        "role": "assistant", "content": answer
    })
    st.rerun()